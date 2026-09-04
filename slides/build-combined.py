#!/usr/bin/env python
"""Merge the three per-day decks into one combined 3-day deck, with a
course-level title cover (not a per-day title).

Run AFTER build-day1/2/3.py have produced their .pptx files.

    cd slides && python build-combined.py   # writes kursus-powerbi-kkdw.pptx
"""
import copy
from io import BytesIO
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE

HERE = Path(__file__).resolve().parent
PARTS = ["day1-fondasi-data.pptx", "day2-power-bi.pptx", "day3-analitik-ai.pptx"]
OUT = HERE / "kursus-powerbi-kkdw.pptx"

BG0   = RGBColor(0x14, 0x17, 0x1F)
GOLD  = RGBColor(0xF2, 0xC8, 0x11)
INK   = RGBColor(0xEE, 0xF1, 0xF6)
MUTED = RGBColor(0x9A, 0xA4, 0xB2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGOLD = RGBColor(0xFF, 0xD8, 0x4D)

COURSE_TITLE = "Visualisasi Data & Dashboard Pintar Berasaskan AI"
COURSE_SUB   = "Power BI · Microsoft Fabric · Copilot"
COURSE_THEME = "Tema: Dashboard Pintar Pemantauan Prestasi Program JPD & BELB Bersepadu dengan MyProjek"
COURSE_FOOT  = "Kementerian Kemajuan Desa dan Wilayah (KKDW) · 3 Hari"


def append_slides(dst, src_path):
    src = Presentation(str(src_path))
    blank = dst.slide_layouts[6]
    for slide in src.slides:
        new = dst.slides.add_slide(blank)
        for shp in list(new.shapes):
            shp._element.getparent().remove(shp._element)
        for shp in slide.shapes:
            # Pictures carry a media part + relationship that a raw XML copy would
            # break, so re-embed them via add_picture (keeps position/size).
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                new.shapes.add_picture(BytesIO(shp.image.blob),
                                       shp.left, shp.top, shp.width, shp.height)
            else:
                new.shapes._spTree.append(copy.deepcopy(shp._element))


def _rect(slide, x, y, w, h, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def add_course_cover(prs):
    """Add a course-level title slide and move it to the front."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, prs.slide_width, prs.slide_height, BG0)
    _rect(s, 0, 0, Inches(0.18), prs.slide_height, GOLD)

    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.4), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = COURSE_TITLE
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Arial"

    tb2 = s.shapes.add_textbox(Inches(0.9), Inches(4.35), Inches(11.4), Inches(2.0))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; r2 = p2.add_run(); r2.text = COURSE_SUB
    r2.font.size = Pt(22); r2.font.bold = True; r2.font.color.rgb = GOLD; r2.font.name = "Arial"
    pt = tf2.add_paragraph(); pt.space_before = Pt(8); rt = pt.add_run(); rt.text = COURSE_THEME
    rt.font.size = Pt(13); rt.font.italic = True; rt.font.color.rgb = LGOLD; rt.font.name = "Arial"
    p3 = tf2.add_paragraph(); p3.space_before = Pt(6); r3 = p3.add_run(); r3.text = COURSE_FOOT
    r3.font.size = Pt(15); r3.font.color.rgb = MUTED; r3.font.name = "Arial"

    # gold accent bar at the bottom
    _rect(s, 0, prs.slide_height - Inches(0.09), prs.slide_width, Inches(0.09), GOLD)

    # move the new (last) slide to the front
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[-1]); lst.insert(0, ids[-1])


missing = [p for p in PARTS if not (HERE / p).exists()]
if missing:
    raise SystemExit(f"Missing decks: {missing}. Run build-day1/2/3.py first.")

combined = Presentation(str(HERE / PARTS[0]))   # start from day1 (keeps size/layout)
for part in PARTS[1:]:                           # append day2 & day3
    append_slides(combined, HERE / part)
add_course_cover(combined)                       # course title, moved to front
combined.core_properties.title = COURSE_TITLE    # document title (not a training/day label)
combined.core_properties.subject = COURSE_SUB
combined.save(str(OUT))
print(f"Wrote {OUT.name} ({len(combined.slides._sldIdLst)} slides)")
