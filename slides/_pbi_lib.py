"""Shared PPTX helpers + Power BI dark/gold theme for the KKDW dashboard course decks.

Used by build-day1.py, build-day2.py, build-day3.py (+ build-combined.py).
Each build script runs in its own process, so the module-level `prs` is fresh per deck.

    import _pbi_lib as L
    s = L.new_slide("Kicker")
    L.title(s, [[("Tajuk ", {}), ("berwarna", {"color": L.GOLD})]])
    ...
    L.prs.save("out.pptx")

Course: BI-FABRIC-KKDW-101 — Visualisasi Data & Dashboard Pintar Berasaskan AI
(Power BI · Microsoft Fabric · Copilot) for KKDW.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- Power BI dark/gold theme ------------------------------------------------
BG0   = RGBColor(0x14, 0x17, 0x1F)   # deep slate
BG1   = RGBColor(0x1C, 0x20, 0x2B)
CARD  = RGBColor(0x26, 0x2B, 0x38)
CARDD = RGBColor(0x33, 0x2C, 0x16)   # gold-tinted "dynamic" card
BRD   = RGBColor(0x3C, 0x44, 0x55)
GOLD  = RGBColor(0xF2, 0xC8, 0x11)   # Power BI yellow — primary accent
LGOLD = RGBColor(0xFF, 0xD8, 0x4D)   # lighter gold (kicker/accents)
BLUE  = RGBColor(0x4A, 0xB3, 0xE0)   # secondary accent
INK   = RGBColor(0xEE, 0xF1, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9A, 0xA4, 0xB2)
GOOD  = RGBColor(0x3D, 0xDC, 0x97)
CODEFG= RGBColor(0xFF, 0xE9, 0xA8)
AMBER = RGBColor(0xF5, 0xC5, 0x6B)

# Back-compat aliases so slide code can use familiar names
PINK  = GOLD     # primary accent (bar, markers, arrows)
VIOLET= LGOLD    # kicker/title accent
PURPLE= BLUE     # numbered circles

SANS = "Arial"
MONO = "Consolas"

SW, SH = 13.333, 7.5
ML, MR, MT = 0.75, 0.75, 0.62
CW = SW - ML - MR

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

TOTAL = 14   # override per deck
COURSE = "Kursus Dashboard Pintar KKDW · Power BI · Fabric · Copilot"


# --- primitives --------------------------------------------------------------
def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background(); shape.shadow.inherit = False

def bg(slide, color=BG0):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _fill(r, color); return r

def accent_bar(slide):
    h = Inches(0.09)
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - h, prs.slide_width, h)
    _fill(r, GOLD)

def box(slide, x, y, w, h, fill=None, line=None, line_w=1.25, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def _apply_runs(p, segments, base_color, base_size, bold_all=False, font=SANS):
    for seg in segments:
        t, o = seg if isinstance(seg, tuple) else (seg, {})
        r = p.add_run(); r.text = t; f = r.font
        f.name = o.get("font", font); f.size = Pt(o.get("size", base_size))
        f.bold = o.get("bold", bold_all); f.italic = o.get("italic", False)
        f.color.rgb = o.get("color", base_color)

def text(slide, x, y, w, h, segments, size=18, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=SANS, line_spacing=1.12, space_after=4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    paras = segments if (isinstance(segments, list) and segments and isinstance(segments[0], list)) else [segments]
    for pi, para in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing; p.space_after = Pt(space_after); p.space_before = Pt(0)
        runs = para if isinstance(para, list) else [(para, {})]
        _apply_runs(p, runs, color, size, bold_all=bold, font=font)
        if italic:
            for r in p.runs: r.font.italic = True
    return tb

def kicker(slide, label):
    k = box(slide, ML, MT, min(len(label)*0.108 + 0.7, 9.5), 0.42, fill=CARD, line=BRD, line_w=1)
    tf = k.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = label.upper()
    r.font.name = SANS; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = LGOLD
    return k

def title(slide, runs, y=1.18, size=38):
    text(slide, ML, y, CW, 1.4, runs, size=size, color=WHITE, bold=True, line_spacing=1.02)

def notes(slide, s):
    slide.notes_slide.notes_text_frame.text = s

def bullets(slide, x, y, w, h, items, size=15.5, color=INK, gap=6, marker="•  "):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.08; p.space_after = Pt(gap); p.space_before = Pt(0)
        mr = p.add_run(); mr.text = marker; mr.font.name = SANS; mr.font.size = Pt(size)
        mr.font.color.rgb = GOLD; mr.font.bold = True
        runs = it if isinstance(it, list) else [(it, {})]
        _apply_runs(p, runs, color, size)
    return tb

def card(slide, x, y, w, h, ico, head, body, dyn=False, head_color=WHITE, body_size=12.5):
    box(slide, x, y, w, h, fill=(CARDD if dyn else CARD), line=(GOLD if dyn else BRD), line_w=1.25)
    pad = 0.22
    if ico: text(slide, x+pad, y+pad-0.02, w-2*pad, 0.5, ico, size=22)
    hy = y + (0.62 if ico else pad)
    text(slide, x+pad, hy, w-2*pad, 0.5, head, size=15, color=head_color, bold=True, line_spacing=1.0)
    if body:
        by = hy + 0.42
        if isinstance(body, list):
            bullets(slide, x+pad, by, w-2*pad, h-(by-y)-pad, body, size=11.5, gap=3, marker="•  ")
        else:
            text(slide, x+pad, by, w-2*pad, h-(by-y)-pad, body, size=body_size, color=MUTED, line_spacing=1.08)

def picture(slide, path, x, y, w, h=None, line=BRD, line_w=1.0):
    """Embed an image (JPG/PNG). Width fixed; height auto-scales if h is None (keeps aspect)."""
    pic = slide.shapes.add_picture(
        path, Inches(x), Inches(y), width=Inches(w),
        height=(Inches(h) if h is not None else None))
    if line is not None:
        pic.line.color.rgb = line; pic.line.width = Pt(line_w)
    pic.shadow.inherit = False
    return pic

def code(slide, x, y, w, h, lines, size=13, title_txt=None):
    """Monospace code/DAX block on a dark card with a gold left rule."""
    box(slide, x, y, w, h, fill=RGBColor(0x0F, 0x12, 0x18), line=BRD, line_w=1)
    box(slide, x, y, 0.06, h, fill=GOLD, radius=False)
    ty = y + 0.14
    if title_txt:
        text(slide, x+0.28, ty, w-0.5, 0.3, title_txt, size=11, color=LGOLD, bold=True, font=MONO)
        ty += 0.34
    segs = [[(ln, {})] for ln in lines]
    text(slide, x+0.28, ty, w-0.5, h-(ty-y)-0.12, segs, size=size, color=CODEFG,
         font=MONO, line_spacing=1.2, space_after=2)

def pipeline(slide, x, y, w, steps, chip_h=0.62, size=13):
    gap = 0.12; arrow_w = 0.34; px = x; py = y; max_x = x + w
    for i, st in enumerate(steps):
        cw = min(0.20 + len(st)*0.098, w)
        if px + cw > max_x and i > 0:
            px = x; py += chip_h + 0.22
        c = box(slide, px, py, cw, chip_h, fill=CARD, line=BRD, line_w=1)
        tf = c.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = st; r.font.name = SANS; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = INK
        px += cw
        if i < len(steps) - 1:
            a = slide.shapes.add_textbox(Inches(px), Inches(py), Inches(arrow_w), Inches(chip_h))
            atf = a.text_frame; atf.vertical_anchor = MSO_ANCHOR.MIDDLE; atf.word_wrap = False
            ap = atf.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
            ar = ap.add_run(); ar.text = "→"; ar.font.name = SANS; ar.font.size = Pt(20); ar.font.bold = True; ar.font.color.rgb = GOLD
            px += arrow_w + gap
    return py + chip_h

def note_strip(slide, s, y=None):
    y = y if y is not None else SH - 1.15
    box(slide, ML, y, 0.04, 0.62, fill=GOLD)
    text(slide, ML+0.22, y-0.02, CW-0.3, 0.75, s, size=12.5, color=MUTED, line_spacing=1.12)

def footer(slide, idx, course=None):
    text(slide, ML, SH-0.42, 9, 0.3, course or COURSE, size=9, color=MUTED)
    text(slide, SW-1.9, SH-0.42, 1.2, 0.3, f"{idx} / {TOTAL}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def new_slide(kicker_label=None):
    s = prs.slides.add_slide(BLANK); bg(s)
    if kicker_label: kicker(s, kicker_label)
    return s

def numbered(slide, x, y, n, color=BLUE, d=0.5):
    b = box(slide, x, y, d, d, fill=color, radius=True)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); r.font.bold = True; r.font.size = Pt(18); r.font.color.rgb = WHITE; r.font.name = SANS
    return b

def table(slide, x, y, w, rows, col_w, header=None, row_h=0.5, size=12.5, head_size=11):
    """rows: list of tuples (str or run-list). col_w: fractional widths ~sum 1."""
    widths = [w*f for f in col_w]; cy = y
    if header:
        cx = x
        for j, htxt in enumerate(header):
            text(slide, cx, cy, widths[j]-0.1, 0.32, htxt, size=head_size, color=LGOLD, bold=True)
            cx += widths[j]
        cy += 0.38
        box(slide, x, cy-0.04, w, 0.02, fill=BRD, radius=False)
    for r in rows:
        cx = x
        for j, cell in enumerate(r):
            text(slide, cx, cy, widths[j]-0.1, row_h, cell, size=size, color=INK, line_spacing=1.02)
            cx += widths[j]
        cy += row_h
    return cy
